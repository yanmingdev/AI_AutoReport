# =============================================================================
# report_generator.py
# =============================================================================

import os
import io
import re
import logging
from datetime import datetime
from pathlib import Path

from dotenv import load_dotenv
import streamlit as st
import streamlit.components.v1 as components
from google import genai
from google.genai import types
from pptx import Presentation
from pptx.util import Pt, Inches  # 🔧 之前已新增：Inches 用來設定 4:3 / 16:9 尺寸
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# =============================================================================
# 0. 基底路徑
# =============================================================================
BASE_DIR = Path(__file__).parent.resolve()

# =============================================================================
# 1. Logging 設定
# =============================================================================
log_dir = str(BASE_DIR / "logs")
os.makedirs(log_dir, exist_ok=True)
today = datetime.now().strftime("%m%d")
log_file = os.path.join(log_dir, f"log_{today}.log")

handlers = [logging.StreamHandler()]
try:
    handlers.insert(0, logging.FileHandler(log_file, encoding="utf-8"))
except Exception:
    pass

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=handlers,
)
logger = logging.getLogger(__name__)
logger.info("=== Application start ===")

# =============================================================================
# 2. 讀取 Gemini API Key
# =============================================================================
load_dotenv(BASE_DIR / ".env")

api_key = os.getenv("GEMINI_API_KEY")
if not api_key:
    try:
        api_key = st.secrets["GEMINI_API_KEY"]
    except Exception:
        api_key = None

if not api_key:
    st.error(
        "❌ 找不到 GEMINI_API_KEY。\n\n"
        "請在 D:\\AI_AutoReport\\.env 內加入：\nGEMINI_API_KEY=你的API金鑰\n\n"
        "（或建立 .streamlit\\secrets.toml 後再放入相同鍵值）"
    )
    st.stop()

# =============================================================================
# 3. Streamlit 頁面設定
# =============================================================================
st.set_page_config(page_title="AI 需求與結案報告產生器", page_icon="✨", layout="wide")

# =============================================================================
# 4. 全域 CSS
# =============================================================================
st.markdown(
    f"""
<style>
:root {{
    --primary-color: {("#FF8C00" if st.session_state.get('doc_type', '結案報告') == "結案報告" else "#1E90FF")};
    --primary-light: {("#FF8C0333" if st.session_state.get('doc_type', '結案報告') == "結案報告" else "#1E90FF33")};
}}
.stButton>button {{
    background-color: var(--primary-color)!important;
    color: #fff!important;
    border: none!important;
}}
[data-testid="stSidebar"] [data-baseweb="tag"] {{ background-color: var(--primary-color)!important; }}
[data-testid="stSidebar"] [data-baseweb="tag"] span,
[data-testid="stSidebar"] [data-baseweb="tag"] svg,
[data-testid="stSidebar"] [data-baseweb="tag-close-button"] {{ color: #fff!important; }}
.rc-slider-rail {{ background-color: var(--primary-light)!important; }}
.rc-slider-track {{ background-color: var(--primary-color)!important; }}
.rc-slider-handle {{ background-color: var(--primary-color)!important; border-color: var(--primary-color)!important; }}
.rc-slider-tooltip-inner {{ background-color: var(--primary-color)!important; border: 1px solid var(--primary-color)!important; color: #fff!important; }}
.rc-slider-handle::after {{ color: #fff!important; }}
.header {{ margin-top:-2.4rem!important; margin-bottom:0!important; padding:0!important; }}
.big-title {{ font-size:28px!important; font-weight:800!important; margin:0; }}
.subtitle {{ font-size:16px!important; color:#ddd!important; margin:0; }}
section[data-testid="stSidebar"] {{ width:260px!important; }}
.block-title {{ font-size:20px!important; margin-top:1rem!important; margin-bottom:0.3rem!important; }}
.stTextArea textarea {{ height:200px!important; }}
[data-testid="stMarkdownContainer"] h4 {{color: var(--primary-color) !important;}}
[data-testid="stExpander"] {{
    border: 1px solid var(--primary-light);
    border-radius: 12px;
    background: rgba(255,255,255,0.02);
    margin-top: 8px;
}}
.stMarkdown hr{{ margin: 0 !important; }}
</style>
""",
    unsafe_allow_html=True,
)

# =============================================================================
# 5. Sidebar
# =============================================================================
st.sidebar.markdown("<p>生成報告格式：</p>", unsafe_allow_html=True)
doc_type = st.sidebar.selectbox(
    "", ["結案報告", "需求文件"], index=0, label_visibility="collapsed"
)
st.session_state["doc_type"] = doc_type

st.sidebar.markdown("<p>目標系統：</p>", unsafe_allow_html=True)
domain_options = ["Generic", "PLM", "SAP/ERP", "Salesforce", "HR", "B2B"]
domain_hint_display = st.sidebar.selectbox(
    "", domain_options, index=0, label_visibility="collapsed"
)

st.sidebar.markdown("<p>選擇要生成的內容區塊：</p>", unsafe_allow_html=True)
available_blocks = [
    "專案名稱",
    "專案目標",
    "專案效益",
    "開發流程",
    "作業時程",
    "專案分工",
]
selected_blocks = st.sidebar.multiselect(
    "區塊", available_blocks, default=[], label_visibility="collapsed"
)

st.sidebar.markdown(
    "<p>創意溫度<br>(0.0＝保守 ↔ 1.0＝創意)</p>", unsafe_allow_html=True
)
creativity_temp = st.sidebar.slider("", 0.0, 1.0, 0.5, 0.1)

# =============================================================================
# 6. 主標題與說明
# =============================================================================
st.markdown(
    f"""
<div class="header">
  <div class="big-title">🚀 AI {doc_type} 產生器</div>
  <div class="subtitle">只要簡單輸入口語化內容，AI 幫你生成專業 {doc_type}！</div>
</div>
""",
    unsafe_allow_html=True,
)

with st.expander("📖 **使用說明**（點我展開）", expanded=False):
    st.markdown(
        f"""
1. 左側必選擇「生成報告格式」與「生成的內容區塊」，建議至少兩個以上。  
2. 可選填「目標系統」。  
3. 可調整「創意溫度」。  
4. 產生後在下方預覽，並可下載 **PPT（4:3 / 16:9）** 與 **Word**。
"""
    )


# =============================================================================
# 7. 載入模板
# =============================================================================
def load_template(path: str) -> str:
    return Path(path).read_text(encoding="utf-8")


# =============================================================================
# 8. 呼叫 Gemini
# =============================================================================
def generate_content(
    project_title: str,
    project_objective: str,
    project_benefit: str,
    development_process: str,
    timeline_schedule: str,
    project_assignment: str,
    domain_hint_display: str,
) -> str:
    template_file = (
        "prompt_template.txt" if doc_type == "結案報告" else "requirement_template.txt"
    )
    tpl_path = BASE_DIR / template_file
    if not tpl_path.exists():
        st.error(f"找不到範本：{tpl_path}")
        st.stop()

    domain_instructions = (
        f"【目標系統：{domain_hint_display}】\n"
        f"- 若為 Generic，請使用通用術語；若為 PLM／SAP/ERP／Salesforce／HR／B2B，請套用對應領域名詞、流程與 KPI。\n\n"
    )

    prompt_body = load_template(str(tpl_path)).format(
        title=project_title,
        goal=project_objective,
        benefit=project_benefit,
        process=development_process,
        schedule=timeline_schedule,
        assignment=project_assignment,
        domain_hint=domain_hint_display,
        domain_hint_display=domain_hint_display,
    )
    prompt = domain_instructions + prompt_body

    client = genai.Client(api_key=api_key)
    cfg = types.GenerateContentConfig(temperature=creativity_temp)
    resp = client.models.generate_content(
        model="gemini-2.5-flash", contents=[prompt], config=cfg
    )
    return resp.text


# =============================================================================
# 9. 動態輸入區
# =============================================================================
field_labels = {
    "專案名稱": "🧩 專案名稱",
    "專案目標": "🎯 專案目標",
    "專案效益": "✨ 專案效益",
    "開發流程": "🛠️ 開發流程",
    "作業時程": "⏳ 作業時程",
    "專案分工": "👥 專案分工",
}
field_values = {}
for i in range(0, len(selected_blocks), 3):
    cols = st.columns(3)
    for j, block in enumerate(selected_blocks[i : i + 3]):
        with cols[j]:
            st.markdown(
                f"<div class='block-title'>{field_labels[block]}</div>",
                unsafe_allow_html=True,
            )
            field_values[block] = st.text_area(
                f"請填寫 {block}：", height=200, label_visibility="collapsed"
            )

project_title = field_values.get("專案名稱", "")
project_objective = field_values.get("專案目標", "")
project_benefit = field_values.get("專案效益", "")
development_process = field_values.get("開發流程", "")
timeline_schedule = field_values.get("作業時程", "")
project_assignment = field_values.get("專案分工", "")

st.write("---")


# =============================================================================
# 10. 從 AI 產文擷取「專案名稱」做檔名
# =============================================================================
def extract_project_title(text):
    patterns = [
        r"一、專案名稱[^\n\r]*\n\s*[-＊*]\s*(.+)",
        r"一、專案名稱[^\n\r]*\n\s*(.+)",
    ]
    for pat in patterns:
        match = re.search(pat, text)
        if match:
            title = match.group(1).strip()
            title = re.sub(r"\s*\(.*?\)$", "", title)
            return title
    return None


# =============================================================================
# 11. session_state
# =============================================================================
if "generated_text" not in st.session_state:
    st.session_state["generated_text"] = ""

# =============================================================================
# 12. 生成按鈕
# =============================================================================
if st.button(f"🪄 生成 {doc_type}", use_container_width=True):
    if not selected_blocks:
        st.warning("請至少選擇一個區塊")
    else:
        missing = [b for b in selected_blocks if not field_values.get(b, "").strip()]
        if missing:
            st.warning("⚠️ 尚未填寫：" + "、".join(missing))
        else:
            with st.spinner("AI 撰寫中..."):
                try:
                    generated_text = generate_content(
                        project_title,
                        project_objective,
                        project_benefit,
                        development_process,
                        timeline_schedule,
                        project_assignment,
                        domain_hint_display,
                    )
                except Exception as e:
                    st.error(f"❌ 發生錯誤：{e}")
                    generated_text = None
            if generated_text:
                st.session_state["generated_text"] = generated_text

# =============================================================================
# 13. 顯示結果 + 下載
# =============================================================================
if st.session_state.get("generated_text"):
    st.success(f"🎉 {doc_type} 生成完成！")
    st.markdown(f"### 📌 {doc_type} 預覽")

    content = st.session_state["generated_text"].strip()
    st.markdown(content)
    st.code(content, language="markdown")

    components.html(
        """
<script>
;(function(){
  try {
    const f = window.frameElement;
    if (f) { f.style.height='0'; f.style.border='0'; f.style.margin='0'; f.style.padding='0'; f.style.minHeight='0'; }
  } catch(e){}
  const bind = ()=>{
    const btn = window.parent.document.querySelector('button[title="Copy to clipboard"]');
    if(!btn||btn.dataset.bound) return;
    btn.dataset.bound = '1';
    const lbl = document.createElement('span');
    lbl.id = 'copy-label';
    lbl.innerText = '點擊複製';
    lbl.style = 'margin-left:8px; color:var(--primary-color); font-size:16px; vertical-align:middle; top:8px; position:relative;';
    btn.parentElement.appendChild(lbl);
    btn.addEventListener('click', ()=>{ lbl.innerText = '已複製'; });
  };
  setInterval(bind,500);
})();
</script>""",
        height=0,
    )

    filename_base = extract_project_title(st.session_state["generated_text"])
    if not filename_base:
        st.error(
            "❌ 無法擷取『專案名稱』（請確認AI回應有『一、專案名稱』區塊），無法下載檔案"
        )
    else:
        filename_base = re.sub(r'[\\/:*?"<>|]', "_", filename_base)
        generated_text = st.session_state["generated_text"]

        # ---------------------------------------------------------------------
        # 🔧 修正：更健壯的標題 regex（允許前綴 - / * + 空白）
        # ---------------------------------------------------------------------
        HEADER_RE = re.compile(r"^\s*(?:[-*]\s*)?(#{1,6})\s+(.+)$", re.M)

        # ---------------------------------------------------------------------
        # 🔧 修正：在 build_pptx() 內「獨立解析」 headers，避免外部狀態干擾
        # ---------------------------------------------------------------------
        def build_pptx(
            markdown_text: str, title_text: str, aspect: str = "4:3"
        ) -> io.BytesIO:
            """
            依據 aspect 建立不同長寬比的 PPT：
              - aspect="4:3"  -> 10in x 7.5in
              - aspect="16:9" -> 13.333in x 7.5in
            回傳 BytesIO 以供 st.download_button 使用。
            """
            ppt = Presentation()

            # 設定投影片尺寸
            if aspect == "16:9":
                ppt.slide_width = Inches(13.3333333)
                ppt.slide_height = Inches(7.5)
            else:  # "4:3"
                ppt.slide_width = Inches(10)
                ppt.slide_height = Inches(7.5)

            # 解析標題（允許 * ## 或 - ## 的行）
            headers_local = list(HEADER_RE.finditer(markdown_text))

            # 首頁
            title_slide_layout = ppt.slide_layouts[0]
            slide = ppt.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = title_text
            if slide.placeholders and len(slide.placeholders) > 1:
                slide.placeholders[1].text = ""

            title_shape = slide.shapes.title
            title_shape.text_frame.paragraphs[0].font.size = Pt(48)
            title_shape.text_frame.paragraphs[0].font.name = "微軟正黑體"
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            def add_slide(title, content):
                s = ppt.slides.add_slide(ppt.slide_layouts[1])
                tf = s.shapes.title.text_frame
                tf.clear()
                tf.margin_top = Pt(5)
                tf.vertical_anchor = MSO_ANCHOR.TOP
                p = tf.paragraphs[0]
                p.text = title
                p.font.name = "微軟正黑體"
                p.font.size = Pt(32)
                p.font.color.rgb = RGBColor(0, 108, 184)
                p.alignment = PP_ALIGN.LEFT

                body = s.placeholders[1].text_frame
                body.clear()
                body.margin_top = Pt(5)
                body.vertical_anchor = MSO_ANCHOR.TOP
                # 行內容：直接逐行加入段落（保留你的原有呈現風格）
                for line in content.split("\n"):
                    para = body.add_paragraph()
                    para.text = line
                    para.font.name = "微軟正黑體"
                    para.font.size = Pt(24)
                    para.font.color.rgb = RGBColor(0, 0, 0)
                    para.alignment = PP_ALIGN.LEFT
                try:
                    body.fit_text(max_size=24)
                except Exception:
                    pass

            if headers_local:
                for idx, h in enumerate(headers_local):
                    start = h.end()
                    end = (
                        headers_local[idx + 1].start()
                        if idx + 1 < len(headers_local)
                        else len(markdown_text)
                    )
                    add_slide(h.group(2).strip(), markdown_text[start:end].strip())
            else:
                # 若真的完全沒有標題，就整段塞一頁
                add_slide(title_text, markdown_text)

            buf = io.BytesIO()
            ppt.save(buf)
            buf.seek(0)
            return buf

        # === 兩種比例的 PPT 下載 ===
        try:
            col_p43, col_p169 = st.columns(2)
            with col_p43:
                buf_43 = build_pptx(generated_text, filename_base, aspect="4:3")
                st.download_button(
                    label="📥 下載 PPT（4:3）",
                    data=buf_43,
                    file_name=f"{filename_base}_4x3.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                )
            with col_p169:
                buf_169 = build_pptx(generated_text, filename_base, aspect="16:9")
                st.download_button(
                    label="📥 下載 PPT（16:9）",
                    data=buf_169,
                    file_name=f"{filename_base}_16x9.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                )
        except ImportError:
            st.error("❌ 無法匯出 PPTX，請 pip install python-pptx")

        # === Word 下載（原樣保留） ===
        try:
            from docx import Document
            from docx.shared import Pt as DocPt

            doc = Document()
            doc.styles["Normal"].font.name = "微軟正黑體"
            doc.styles["Normal"].font.size = DocPt(12)

            headers_for_doc = list(HEADER_RE.finditer(generated_text))
            if headers_for_doc:
                for idx, h in enumerate(headers_for_doc):
                    start = h.end()
                    end = (
                        headers_for_doc[idx + 1].start()
                        if idx + 1 < len(headers_for_doc)
                        else len(generated_text)
                    )
                    heading = h.group(2).strip()
                    lines = generated_text[start:end].strip().split("\n")
                    doc.add_heading(heading, level=2)
                    for line in lines:
                        if line.strip():
                            p = doc.add_paragraph(line)
                            p.style = doc.styles["Normal"]
            else:
                doc.add_paragraph(generated_text)

            doc_buf = io.BytesIO()
            doc.save(doc_buf)
            doc_buf.seek(0)
            st.download_button(
                label="📥 下載 Word 檔",
                data=doc_buf,
                file_name=f"{filename_base}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True,
            )
        except ImportError:
            st.error("❌ 無法匯出 Word，請 pip install python-docx")
